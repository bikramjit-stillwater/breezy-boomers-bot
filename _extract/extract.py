"""
extract.py — Pull 100% of the Urban Hipster data from the two source files:
  1. Excel: 'Persona Info' sheet (all profile/metric fields) + 'Urban Hipster'
     sheet (the full brand spend-propensity matrix).
  2. PPTX: all text from slides that mention Urban Hipster.

Outputs _extract/urban_hipster_source.json  (raw, complete capture).
"""
import json, os, re
import openpyxl
from pptx import Presentation

HERE = os.path.dirname(__file__)
ROOT = os.path.dirname(HERE)
XLSX = os.path.join(ROOT, "Fremantle 2025 Segmentation AI-Ready Format 2026_05 V2.xlsx")
PPTX = os.path.join(ROOT, "Fremantle_Personas_050426_V1.2.pptx")
SEG = "Urban Hipster"
# PPTX slides label this segment in both singular and plural ("Urban Hipster(s)").
PPTX_RE = r"urban\s*hipster"

out = {"segment": SEG, "excel_persona_info": {}, "spend_propensity": {}, "pptx_text": []}

# ---------------------------------------------------------------------------
# 1. Persona Info sheet -> all fields for Urban Hipster
# ---------------------------------------------------------------------------
wb = openpyxl.load_workbook(XLSX, data_only=True)
pi = wb["Persona Info"]
headers = [pi.cell(1, c).value for c in range(1, pi.max_column + 1)]
# Find the Urban Hipster row
bb_row = None
for r in range(2, pi.max_row + 1):
    if str(pi.cell(r, 1).value).strip() == SEG:
        bb_row = r
        break
fields = {}
for c, h in enumerate(headers, start=1):
    if h is None:
        continue
    v = pi.cell(bb_row, c).value
    if v is not None and str(v).strip() != "":
        fields[str(h).strip()] = v
out["excel_persona_info"] = fields
print(f"Persona Info: captured {len(fields)} fields for {SEG}")

# ---------------------------------------------------------------------------
# 2. Urban Hipster sheet -> 6-block spend-propensity matrix
#    Blocks start at columns: B(2), H(8), N(14), T(20), Z(26), AF(32)
#    Each block: [name, score, rounded, index] with a 'Deck Order' col before.
# ---------------------------------------------------------------------------
bs = wb[SEG]
block_name_cols = {}  # category title -> name column index
for c in range(1, bs.max_column + 1):
    title = bs.cell(1, c).value
    if title and str(title).strip() not in ("Deck Order", "Score", "Rounded", "Index"):
        block_name_cols[str(title).strip()] = c

spend = {}
total_brands = 0
for category, name_col in block_name_cols.items():
    idx_col = name_col + 3  # name, score, rounded, index
    brands = {}
    for r in range(2, bs.max_row + 1):
        brand = bs.cell(r, name_col).value
        index = bs.cell(r, idx_col).value
        if brand is not None and str(brand).strip() != "":
            brands[str(brand).strip()] = (str(index).strip() if index is not None else None)
    if brands:
        spend[category] = brands
        total_brands += len(brands)
out["spend_propensity"] = spend
print(f"Spend propensity: {len(spend)} categories, {total_brands} brands total")

# ---------------------------------------------------------------------------
# 3. PPTX -> text + tables + CHARTS from slides mentioning Urban Hipster.
#    Recurses into grouped shapes and reads native chart data (category:value),
#    which a plain text-frame extractor would otherwise miss.
# ---------------------------------------------------------------------------
from pptx.enum.shapes import MSO_SHAPE_TYPE
prs = Presentation(PPTX)
out["pptx_charts"] = []  # [{slide, title, points:[{label,value}]}]

def shape_text(shape, charts_sink, slide_no):
    chunks = []
    # Recurse into groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            chunks.extend(shape_text(sub, charts_sink, slide_no))
        return chunks
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs)
            if t.strip():
                chunks.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text for c in row.cells]
            if any(x.strip() for x in cells):
                chunks.append(" | ".join(cells))
    if shape.has_chart:
        try:
            ch = shape.chart
            cats = [str(c) for c in ch.plots[0].categories]
            title = ch.chart_title.text_frame.text if ch.has_title else ""
            for series in ch.series:
                pts = [{"label": cats[j] if j < len(cats) else str(j),
                        "value": (round(v, 1) if isinstance(v, float) else v)}
                       for j, v in enumerate(series.values)]
                charts_sink.append({"slide": slide_no,
                                    "title": (title or series.name or "chart"),
                                    "points": pts})
        except Exception as e:
            charts_sink.append({"slide": slide_no, "title": "unreadable-chart", "error": str(e)})
    return chunks

for i, slide in enumerate(prs.slides, start=1):
    texts = []
    slide_charts = []
    for shape in slide.shapes:
        texts.extend(shape_text(shape, slide_charts, i))
    joined = "\n".join(texts)
    if re.search(PPTX_RE, joined, re.I):
        out["pptx_text"].append({"slide": i, "text": joined})
        out["pptx_charts"].extend(slide_charts)
print(f"PPTX: captured {len(out['pptx_text'])} Urban Hipster slides, "
      f"{len(out['pptx_charts'])} charts")

os.makedirs(HERE, exist_ok=True)
with open(os.path.join(HERE, "urban_hipster_source.json"), "w", encoding="utf-8") as f:
    json.dump(out, f, indent=2, ensure_ascii=False)
print("Wrote _extract/urban_hipster_source.json")
